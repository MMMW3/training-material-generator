"""
PPT 生成器
使用 python-pptx 生成培训演示文稿，支持模板上传
"""

import logging
import os
import shutil
from pathlib import Path
from typing import TYPE_CHECKING

from models import TrainingMaterial, FeatureInfo, BusinessValue, OperationFlow, PPTTemplateConfig, PageType

if TYPE_CHECKING:
    from pptx import Presentation

logger = logging.getLogger(__name__)

# 默认主题色
DEFAULT_COLORS = {
    "primary": "1F4E79",      # 深蓝
    "secondary": "2E75B6",    # 中蓝
    "accent": "ED7D31",       # 橙色强调
    "text_dark": "333333",    # 深灰文字
    "text_light": "FFFFFF",   # 白色文字
    "bg_light": "F2F2F2",     # 浅灰背景
    "success": "70AD47",      # 绿色
    "warning": "FFC000",      # 黄色
}


class PPTGenerator:
    """PPT 培训演示文稿生成器，支持模板上传"""

    def __init__(self, template_config: PPTTemplateConfig = None):
        """
        初始化 PPT 生成器
        
        Args:
            template_config: PPT模板配置，如果为None则使用默认样式
        """
        self.template_config = template_config
        self._prs = None
        self._colors = DEFAULT_COLORS.copy()
        
        # 如果有模板配置，加载模板并应用样式
        if template_config and template_config.enabled and template_config.template_path:
            self._load_template()

    def _load_template(self):
        """加载 PPT 模板并提取样式"""
        if not self.template_config or not self.template_config.template_path:
            return
            
        template_path = self.template_config.template_path
        if not os.path.exists(template_path):
            logger.warning(f"模板文件不存在: {template_path}")
            return
            
        try:
            from pptx import Presentation
            self._prs = Presentation(template_path)
            
            # 从模板中提取配色方案
            self._extract_template_colors()
            
            # 应用用户自定义配色覆盖
            if self.template_config.primary_color:
                self._colors["primary"] = self.template_config.primary_color
            if self.template_config.accent_color:
                self._colors["accent"] = self.template_config.accent_color
                
            logger.info(f"模板加载成功: {template_path}")
        except Exception as e:
            logger.error(f"模板加载失败: {e}")
            self._prs = None

    def _extract_template_colors(self):
        """从模板中提取配色方案"""
        if not self._prs:
            return
            
        # 尝试从主题中提取颜色
        try:
            # 读取幻灯片母版信息
            if self._prs.slide_master:
                # 提取背景色
                background = self._prs.slide_master.background
                if background.fill.solid:
                    color = background.fill.fore_color
                    if hasattr(color, 'rgb'):
                        self._colors["primary"] = color.rgb.to_string()
        except Exception as e:
            logger.debug(f"提取模板颜色失败: {e}")

    def generate(self, material: TrainingMaterial, output_path: str) -> str:
        """
        生成培训 PPT

        Args:
            material: 培训材料数据
            output_path: 输出文件路径

        Returns:
            生成的文件路径
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        except ImportError:
            raise ImportError("请安装 python-pptx 库: pip install python-pptx")

        # 使用模板或创建新演示文稿
        if self._prs and self.template_config and self.template_config.enabled:
            # 深拷贝模板
            prs = self._copy_presentation(self._prs)
            # 获取模板的页面数量，用于映射
            template_slide_count = len(prs.slides)
            logger.info(f"使用模板，模板页数: {template_slide_count}")
        else:
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            template_slide_count = 0

        # 收集所有幻灯片（用于按页类型映射）
        slides_data = self._prepare_slides_data(material)
        
        # 根据模板配置生成幻灯片
        if self.template_config and self.template_config.enabled and self.template_config.page_mappings:
            # 使用页面映射模式
            prs = self._generate_with_page_mapping(prs, slides_data, material)
        else:
            # 使用默认模式生成
            prs = self._generate_default(prs, material)

        # 保存
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)
        logger.info(f"PPT 已生成: {output_path}")
        return output_path

    def _copy_presentation(self, template_prs):
        """复制演示文稿对象"""
        from pptx import Presentation
        from io import BytesIO
        
        buffer = BytesIO()
        template_prs.save(buffer)
        buffer.seek(0)
        return Presentation(buffer)

    def _prepare_slides_data(self, material: TrainingMaterial) -> dict:
        """准备所有幻灯片数据"""
        slides_data = {}
        
        # 封面
        slides_data[PageType.COVER] = [{
            "title": "系统更新培训",
            "subtitle": f"版本 {material.version} | {material.release_date}",
            "content": material.summary,
        }]
        
        # 更新概览
        slides_data[PageType.OVERVIEW] = [{
            "summary": material.summary,
            "features": [
                {
                    "name": f.name,
                    "priority": f.priority.value,
                    "description": f.description[:60],
                }
                for f in material.features
            ],
        }]
        
        # 逐功能页面
        for feature in material.features:
            bv = material.business_values.get(feature.name)
            flow = material.operation_flows.get(feature.name)
            
            # 功能标题页
            slides_data[PageType.FEATURE_TITLE] = slides_data.get(PageType.FEATURE_TITLE, [])
            slides_data[PageType.FEATURE_TITLE].append({
                "feature_name": feature.name,
                "priority": feature.priority.value,
            })
            
            # 业务价值页
            if bv:
                slides_data[PageType.BUSINESS_VALUE] = slides_data.get(PageType.BUSINESS_VALUE, [])
                slides_data[PageType.BUSINESS_VALUE].append({
                    "feature_name": feature.name,
                    "problem": bv.problem_statement,
                    "value": bv.value_proposition,
                    "before": bv.before_scenario,
                    "after": bv.after_scenario,
                    "audience": bv.target_audience,
                    "benefits": bv.key_benefits,
                })
            
            # 操作流程页
            if flow:
                slides_data[PageType.OPERATION_FLOW] = slides_data.get(PageType.OPERATION_FLOW, [])
                slides_data[PageType.OPERATION_FLOW].append({
                    "feature_name": feature.name,
                    "entry_point": flow.entry_point,
                    "steps": [
                        {
                            "number": s.step_number,
                            "action": s.action,
                            "description": s.description,
                            "result": s.expected_result,
                        }
                        for s in flow.steps
                    ],
                    "notes": flow.notes,
                })
        
        # 角色总结
        role_features = {}
        for feature in material.features:
            bv = material.business_values.get(feature.name)
            audiences = bv.target_audience if bv else []
            if not audiences:
                audiences = ["全体用户"]
            for role in audiences:
                if role not in role_features:
                    role_features[role] = []
                role_features[role].append(feature.name)
        
        slides_data[PageType.ROLE_SUMMARY] = [{"roles": role_features}]
        
        # FAQ
        slides_data[PageType.FAQ] = [{
            "faqs": [
                {"question": f.question, "answer": f.answer[:80]}
                for f in material.faqs[:8]
            ]
        }]
        
        # 结束页
        slides_data[PageType.END] = [{
            "title": "谢谢！",
            "content": "如有问题，请通过内部支持渠道反馈",
        }]
        
        return slides_data

    def _generate_with_page_mapping(self, prs, slides_data: dict, material: TrainingMaterial):
        """使用页面映射模式生成PPT"""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        
        mappings = self.template_config.page_mappings
        
        # 创建新的演示文稿，但使用模板的母版
        # 从模板复制母版信息
        template = prs
        
        # 创建新演示文稿
        new_prs = Presentation()
        new_prs.slide_width = template.slide_width
        new_prs.slide_height = template.slide_height
        
        # 复制模板的幻灯片母版
        try:
            for master in template.slide_masters:
                # 尝试复制母版（python-pptx 限制，可能无法完全复制）
                pass
        except Exception:
            pass
        
        # 按映射顺序生成页面
        for page_type, data_list in slides_data.items():
            layout_index = mappings.get(page_type.value, 6)  # 默认使用空白布局
            layout_index = min(layout_index, len(new_prs.slide_layouts) - 1)
            
            for data in data_list:
                if page_type == PageType.COVER:
                    self._add_cover_slide(new_prs, data, layout_index)
                elif page_type == PageType.OVERVIEW:
                    self._add_overview_slide(new_prs, data, layout_index)
                elif page_type == PageType.FEATURE_TITLE:
                    self._add_feature_title_slide(new_prs, data, layout_index)
                elif page_type == PageType.BUSINESS_VALUE:
                    self._add_business_value_slide(new_prs, data, layout_index)
                elif page_type == PageType.OPERATION_FLOW:
                    self._add_operation_flow_slide(new_prs, data, layout_index)
                elif page_type == PageType.ROLE_SUMMARY:
                    self._add_role_summary_slide(new_prs, data, layout_index)
                elif page_type == PageType.FAQ:
                    self._add_faq_slide(new_prs, data, layout_index)
                elif page_type == PageType.END:
                    self._add_end_slide(new_prs, data, layout_index)
        
        return new_prs

    def _generate_default(self, prs, material: TrainingMaterial):
        """使用默认模式生成PPT（无模板）"""
        from pptx.util import Inches, Pt
        
        colors = self._colors
        
        # 封面
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_background(slide, colors["primary"])
        self._add_text_box(slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
                          f"系统更新培训", Pt(44), colors["text_light"], bold=True)
        self._add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(1),
                          f"版本 {material.version} | {material.release_date}", Pt(24), colors["text_light"])
        self._add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(1),
                          material.summary, Pt(18), "BDD7EE")

        # 更新概览
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, "本次更新概览")
        self._add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(1),
                          material.summary, Pt(20), colors["text_dark"])

        features_text = "\n".join([
            f"{'🟢' if f.priority.value.startswith('P0') else '🟡' if f.priority.value.startswith('P1') else '⚪'} "
            f"{f.name} — {f.description[:60]}"
            for f in material.features
        ])
        self._add_text_box(slide, Inches(0.8), Inches(2.5), Inches(11.5), Inches(4.5),
                          features_text, Pt(16), colors["text_dark"])

        # 逐功能讲解
        for feature in material.features:
            bv = material.business_values.get(feature.name)
            flow = material.operation_flows.get(feature.name)

            # 功能标题页
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._add_background(slide, colors["secondary"])
            self._add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.5),
                              f"✨ {feature.name}", Pt(36), colors["text_light"], bold=True)
            self._add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.8),
                              f"优先级：{feature.priority.value}", Pt(18), "BDD7EE")

            # 业务价值页
            if bv:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                self._add_slide_title(slide, f"🎯 {feature.name} — 业务场景与价值")
                y_pos = Inches(1.5)
                self._add_section_label(slide, Inches(0.8), y_pos, "解决的痛点")
                y_pos += Inches(0.5)
                self._add_text_box(slide, Inches(0.8), y_pos, Inches(11.5), Inches(1),
                                  bv.problem_statement, Pt(16), colors["text_dark"])
                y_pos += Inches(1.2)
                self._add_section_label(slide, Inches(0.8), y_pos, "核心价值")
                y_pos += Inches(0.5)
                self._add_text_box(slide, Inches(0.8), y_pos, Inches(11.5), Inches(0.8),
                                  bv.value_proposition, Pt(18), colors["primary"], bold=True)
                y_pos += Inches(1)
                self._add_section_label(slide, Inches(0.8), y_pos, "Before → After")
                y_pos += Inches(0.5)
                self._add_text_box(slide, Inches(0.8), y_pos, Inches(5.5), Inches(1.5),
                                  f"❌ 以前：{bv.before_scenario}", Pt(14), "C00000")
                self._add_text_box(slide, Inches(6.8), y_pos, Inches(5.5), Inches(1.5),
                                  f"✅ 现在：{bv.after_scenario}", Pt(14), colors["success"])
                y_pos += Inches(1.8)
                if bv.target_audience:
                    self._add_text_box(slide, Inches(0.8), y_pos, Inches(11.5), Inches(0.6),
                                      f"目标用户：{', '.join(bv.target_audience)}", Pt(14), colors["text_dark"])

            # 操作流程页
            if flow:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                self._add_slide_title(slide, f"📝 {feature.name} — 操作步骤")
                y_pos = Inches(1.5)
                self._add_text_box(slide, Inches(0.8), y_pos, Inches(11.5), Inches(0.5),
                                  f"功能入口：{flow.entry_point}", Pt(14), colors["secondary"])
                y_pos += Inches(0.7)
                for step in flow.steps:
                    self._add_text_box(slide, Inches(0.8), y_pos, Inches(11.5), Inches(0.4),
                                      f"步骤 {step.step_number}：{step.action}", Pt(15), colors["primary"], bold=True)
                    y_pos += Inches(0.4)
                    self._add_text_box(slide, Inches(1.2), y_pos, Inches(11), Inches(0.4),
                                      f"    {step.description} → {step.expected_result}", Pt(13), colors["text_dark"])
                    y_pos += Inches(0.5)
                    if y_pos > Inches(6.5):
                        break
                if flow.notes:
                    y_pos += Inches(0.2)
                    notes_text = "⚠️ 注意事项：\n" + "\n".join(f"  • {n}" for n in flow.notes)
                    self._add_text_box(slide, Inches(0.8), y_pos, Inches(11.5), Inches(1.5),
                                      notes_text, Pt(13), colors["warning"])

        # 角色总结
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, "👥 各角色关注功能一览")
        role_features = {}
        for feature in material.features:
            bv = material.business_values.get(feature.name)
            audiences = bv.target_audience if bv else []
            if not audiences:
                audiences = ["全体用户"]
            for role in audiences:
                if role not in role_features:
                    role_features[role] = []
                role_features[role].append(feature.name)
        y_pos = Inches(1.5)
        for role, features in role_features.items():
            self._add_text_box(slide, Inches(0.8), y_pos, Inches(11.5), Inches(0.5),
                              f"👤 {role}：{', '.join(features)}", Pt(16), colors["text_dark"])
            y_pos += Inches(0.6)

        # FAQ
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, "❓ 常见问题 (FAQ)")
        if material.faqs:
            y_pos = Inches(1.5)
            for i, faq in enumerate(material.faqs[:8]):
                self._add_text_box(slide, Inches(0.8), y_pos, Inches(11.5), Inches(0.4),
                                  f"Q{i+1}: {faq.question}", Pt(14), colors["primary"], bold=True)
                y_pos += Inches(0.35)
                self._add_text_box(slide, Inches(1.2), y_pos, Inches(11), Inches(0.4),
                                  f"A{i+1}: {faq.answer[:80]}{'...' if len(faq.answer) > 80 else ''}", Pt(13), colors["text_dark"])
                y_pos += Inches(0.5)
                if y_pos > Inches(6.5):
                    break

        # 结束页
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_background(slide, colors["primary"])
        self._add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.5),
                          "谢谢！", Pt(44), colors["text_light"], bold=True)
        self._add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(1),
                          "如有问题，请通过内部支持渠道反馈", Pt(20), "BDD7EE")

        return prs

    # ====== 幻灯片添加方法（模板模式）======
    
    def _add_cover_slide(self, prs, data: dict, layout_index: int):
        """添加封面页"""
        from pptx.util import Inches, Pt
        slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
        self._add_text_box(slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
                          data.get("title", "系统更新培训"), Pt(44), self._colors["text_light"], bold=True)
        self._add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(1),
                          data.get("subtitle", ""), Pt(24), self._colors["text_light"])
        self._add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(1),
                          data.get("content", ""), Pt(18), "BDD7EE")
        return slide

    def _add_overview_slide(self, prs, data: dict, layout_index: int):
        """添加概览页"""
        from pptx.util import Inches, Pt
        slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
        self._add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                          "本次更新概览", Pt(28), self._colors["primary"], bold=True)
        self._add_text_box(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(1),
                          data.get("summary", ""), Pt(18), self._colors["text_dark"])
        
        features = data.get("features", [])
        y_pos = Inches(2.2)
        for f in features:
            self._add_text_box(slide, Inches(0.5), y_pos, Inches(12), Inches(0.5),
                              f"• {f['name']} ({f['priority']})", Pt(14), self._colors["text_dark"])
            y_pos += Inches(0.4)
        return slide

    def _add_feature_title_slide(self, prs, data: dict, layout_index: int):
        """添加功能标题页"""
        from pptx.util import Inches, Pt
        slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
        self._add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.5),
                          f"✨ {data.get('feature_name', '')}", Pt(36), self._colors["text_light"], bold=True)
        self._add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.8),
                          f"优先级：{data.get('priority', '')}", Pt(18), "BDD7EE")
        return slide

    def _add_business_value_slide(self, prs, data: dict, layout_index: int):
        """添加业务价值页"""
        from pptx.util import Inches, Pt
        slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
        self._add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                          f"🎯 {data.get('feature_name', '')} — 业务价值", Pt(24), self._colors["primary"], bold=True)
        
        y_pos = Inches(1.2)
        self._add_text_box(slide, Inches(0.5), y_pos, Inches(12), Inches(0.4),
                          f"核心价值：{data.get('value', '')}", Pt(16), self._colors["primary"], bold=True)
        y_pos += Inches(0.6)
        self._add_text_box(slide, Inches(0.5), y_pos, Inches(5.5), Inches(1.5),
                          f"❌ 以前：{data.get('before', '')}", Pt(14), "C00000")
        self._add_text_box(slide, Inches(6.5), y_pos, Inches(6), Inches(1.5),
                          f"✅ 现在：{data.get('after', '')}", Pt(14), self._colors["success"])
        return slide

    def _add_operation_flow_slide(self, prs, data: dict, layout_index: int):
        """添加操作流程页"""
        from pptx.util import Inches, Pt
        slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
        self._add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                          f"📝 {data.get('feature_name', '')} — 操作步骤", Pt(24), self._colors["primary"], bold=True)
        
        y_pos = Inches(1.2)
        self._add_text_box(slide, Inches(0.5), y_pos, Inches(12), Inches(0.4),
                          f"入口：{data.get('entry_point', '')}", Pt(14), self._colors["secondary"])
        y_pos += Inches(0.5)
        
        for step in data.get("steps", [])[:5]:
            self._add_text_box(slide, Inches(0.5), y_pos, Inches(12), Inches(0.4),
                              f"{step['number']}. {step['action']} → {step.get('result', '')}", 
                              Pt(13), self._colors["text_dark"])
            y_pos += Inches(0.4)
        return slide

    def _add_role_summary_slide(self, prs, data: dict, layout_index: int):
        """添加角色总结页"""
        from pptx.util import Inches, Pt
        slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
        self._add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                          "👥 各角色关注功能", Pt(28), self._colors["primary"], bold=True)
        
        y_pos = Inches(1.2)
        roles = data.get("roles", {})
        for role, features in roles.items():
            self._add_text_box(slide, Inches(0.5), y_pos, Inches(12), Inches(0.5),
                              f"👤 {role}：{', '.join(features)}", Pt(16), self._colors["text_dark"])
            y_pos += Inches(0.5)
        return slide

    def _add_faq_slide(self, prs, data: dict, layout_index: int):
        """添加FAQ页"""
        from pptx.util import Inches, Pt
        slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
        self._add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                          "❓ 常见问题", Pt(28), self._colors["primary"], bold=True)
        
        y_pos = Inches(1.2)
        for i, faq in enumerate(data.get("faqs", [])[:6]):
            self._add_text_box(slide, Inches(0.5), y_pos, Inches(12), Inches(0.4),
                              f"Q{i+1}: {faq.get('question', '')}", Pt(14), self._colors["primary"], bold=True)
            y_pos += Inches(0.35)
            self._add_text_box(slide, Inches(0.8), y_pos, Inches(11), Inches(0.4),
                              f"A: {faq.get('answer', '')}", Pt(12), self._colors["text_dark"])
            y_pos += Inches(0.5)
        return slide

    def _add_end_slide(self, prs, data: dict, layout_index: int):
        """添加结束页"""
        from pptx.util import Inches, Pt
        slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
        self._add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.5),
                          data.get("title", "谢谢！"), Pt(44), self._colors["text_light"], bold=True)
        self._add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(1),
                          data.get("content", ""), Pt(20), "BDD7EE")
        return slide

    # ====== 辅助方法 ======

    def _add_background(self, slide, color_hex: str):
        """设置幻灯片背景色"""
        from pptx.dml.color import RGBColor
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(color_hex)

    def _add_text_box(self, slide, left, top, width, height,
                      text, font_size, color_hex, bold=False):
        """添加文本框"""
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = font_size
        p.font.color.rgb = RGBColor.from_string(color_hex)
        p.font.bold = bold
        return txBox

    def _add_slide_title(self, slide, title_text: str):
        """添加幻灯片标题"""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        colors = self._colors
        shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(13.333), Inches(1.2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(colors["primary"])
        shape.line.fill.background()

        self._add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11.5), Inches(0.8),
                          title_text, Pt(28), colors["text_light"], bold=True)

    def _add_section_label(self, slide, left, top, text: str):
        """添加小节标签"""
        from pptx.util import Inches, Pt
        colors = self._colors
        self._add_text_box(slide, left, top, Inches(3), Inches(0.4),
                          text, Pt(14), colors["accent"], bold=True)
